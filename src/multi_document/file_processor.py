# multi_document/file_processor.py - Multi-Format File Processing
"""
Multi-Format File Processor for RAG System

Supports:
- PDF (existing)
- DOCX (Microsoft Word)
- PPTX (PowerPoint) 
- TXT (Plain text)
- CSV (Structured data)
"""

import os
from typing import List, Dict, Optional
from langchain_core.documents import Document
from langchain_community.document_loaders import PyPDFLoader
from pptx import Presentation
import mimetypes

# DOCX support
try:
    from docx import Document as DocxDocument
    DOCX_AVAILABLE = True
except ImportError:
    DOCX_AVAILABLE = False
    print("⚠️  python-docx not installed. DOCX support disabled.")

# CSV support
try:
    import pandas as pd
    CSV_AVAILABLE = True
except ImportError:
    CSV_AVAILABLE = False
    print("⚠️  pandas not installed. CSV support disabled.")


class FileProcessor:
    """
    Multi-format file processor for RAG system
    """
    
    SUPPORTED_FORMATS = {
        'pdf': 'PDF Document',
        'docx': 'Microsoft Word',
        'pptx': 'PowerPoint Presentation',
        'txt': 'Plain Text',
        'csv': 'CSV Data'
    }
    
    def __init__(self):
        """Initialize file processor"""
        self.format_handlers = {
            'pdf': self.process_pdf,
            'docx': self.process_docx,
            'pptx': self.process_pptx,
            'txt': self.process_txt,
            'csv': self.process_csv
        }
    
    def auto_detect_format(self, file_path: str) -> Optional[str]:
        """
        Auto-detect file format from extension and MIME type
        
        Args:
            file_path: Path to file
            
        Returns:
            Format string (pdf, docx, pptx, txt, csv) or None
        """
        # Get extension
        _, ext = os.path.splitext(file_path)
        ext = ext.lower().lstrip('.')
        
        # Check if supported
        if ext in self.SUPPORTED_FORMATS:
            return ext
        
        # Try MIME type detection
        mime_type, _ = mimetypes.guess_type(file_path)
        if mime_type:
            mime_map = {
                'application/pdf': 'pdf',
                'application/vnd.openxmlformats-officedocument.wordprocessingml.document': 'docx',
                'application/vnd.openxmlformats-officedocument.presentationml.presentation': 'pptx',
                'text/plain': 'txt',
                'text/csv': 'csv'
            }
            return mime_map.get(mime_type)
        
        return None
    
    def process_file(self, file_path: str, metadata: Optional[Dict] = None) -> List[Document]:
        """
        Process file based on format
        
        Args:
            file_path: Path to file
            metadata: Additional metadata
            
        Returns:
            List of Document objects
        """
        # Auto-detect format
        file_format = self.auto_detect_format(file_path)
        
        if not file_format:
            raise ValueError(f"Unsupported file format: {file_path}")
        
        # Get handler
        handler = self.format_handlers.get(file_format)
        if not handler:
            raise ValueError(f"No handler for format: {file_format}")
        
        # Process file
        documents = handler(file_path)
        
        # Add custom metadata
        if metadata:
            for doc in documents:
                doc.metadata.update(metadata)
        
        # Add file format to metadata
        for doc in documents:
            doc.metadata['file_type'] = file_format
            doc.metadata['file_size'] = os.path.getsize(file_path)
        
        return documents
    
    def process_pdf(self, file_path: str) -> List[Document]:
        """
        Process PDF file
        
        Args:
            file_path: Path to PDF
            
        Returns:
            List of Document objects (1 per page)
        """
        loader = PyPDFLoader(file_path)
        documents = loader.load()
        
        print(f"✅ PDF: {len(documents)} pages from {os.path.basename(file_path)}")
        return documents
    
    def process_docx(self, file_path: str) -> List[Document]:
        """
        Process DOCX file (Microsoft Word)
        
        Args:
            file_path: Path to DOCX
            
        Returns:
            List of Document objects
        """
        if not DOCX_AVAILABLE:
            raise ImportError("python-docx not installed. Run: pip install python-docx")
        
        doc = DocxDocument(file_path)
        documents = []
        
        # Extract paragraphs
        paragraphs = []
        for para in doc.paragraphs:
            if para.text.strip():
                paragraphs.append(para.text)
        
        # Extract tables
        tables_text = []
        for table in doc.tables:
            table_data = []
            for row in table.rows:
                row_data = [cell.text.strip() for cell in row.cells]
                table_data.append(" | ".join(row_data))
            tables_text.append("\n".join(table_data))
        
        # Combine content
        content = "\n\n".join(paragraphs)
        if tables_text:
            content += "\n\n" + "\n\n".join(tables_text)
        
        # Create document
        doc_obj = Document(
            page_content=content,
            metadata={
                "source": file_path,
                "filename": os.path.basename(file_path),
                "paragraph_count": len(paragraphs),
                "table_count": len(doc.tables)
            }
        )
        documents.append(doc_obj)
        
        print(f"✅ DOCX: {len(paragraphs)} paragraphs, {len(doc.tables)} tables from {os.path.basename(file_path)}")
        return documents
    
    def process_pptx(self, file_path: str) -> List[Document]:
        """
        Process PPTX file (PowerPoint)
        
        Args:
            file_path: Path to PPTX
            
        Returns:
            List of Document objects (1 per slide)
        """
        prs = Presentation(file_path)
        documents = []
        
        for slide_idx, slide in enumerate(prs.slides):
            text_content = []
            
            # Extract text from shapes
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    text_content.append(shape.text)
            
            # Create document if has content
            if text_content:
                content = "\n".join(text_content)
                doc = Document(
                    page_content=content,
                    metadata={
                        "source": file_path,
                        "filename": os.path.basename(file_path),
                        "slide": slide_idx + 1,
                        "total_slides": len(prs.slides)
                    }
                )
                documents.append(doc)
        
        print(f"✅ PPTX: {len(documents)} slides from {os.path.basename(file_path)}")
        return documents
    
    def process_txt(self, file_path: str) -> List[Document]:
        """
        Process TXT file (Plain text)
        
        Args:
            file_path: Path to TXT
            
        Returns:
            List of Document objects
        """
        # Try different encodings
        encodings = ['utf-8', 'latin-1', 'cp1252']
        content = None
        
        for encoding in encodings:
            try:
                with open(file_path, 'r', encoding=encoding) as f:
                    content = f.read()
                break
            except UnicodeDecodeError:
                continue
        
        if content is None:
            raise ValueError(f"Could not decode file: {file_path}")
        
        # Create document
        doc = Document(
            page_content=content,
            metadata={
                "source": file_path,
                "filename": os.path.basename(file_path),
                "char_count": len(content),
                "line_count": content.count('\n') + 1
            }
        )
        
        print(f"✅ TXT: {len(content)} characters from {os.path.basename(file_path)}")
        return [doc]
    
    def process_csv(self, file_path: str) -> List[Document]:
        """
        Process CSV file (Structured data)
        
        Args:
            file_path: Path to CSV
            
        Returns:
            List of Document objects
        """
        if not CSV_AVAILABLE:
            raise ImportError("pandas not installed. Run: pip install pandas")
        
        # Read CSV
        df = pd.read_csv(file_path)
        
        # Convert to text format
        content_parts = []
        
        # Add header
        content_parts.append(f"CSV File: {os.path.basename(file_path)}")
        content_parts.append(f"Columns: {', '.join(df.columns)}")
        content_parts.append(f"Rows: {len(df)}")
        content_parts.append("")
        
        # Add data (limit to first 100 rows for performance)
        max_rows = min(100, len(df))
        content_parts.append(df.head(max_rows).to_string(index=False))
        
        if len(df) > max_rows:
            content_parts.append(f"\n... ({len(df) - max_rows} more rows)")
        
        content = "\n".join(content_parts)
        
        # Create document
        doc = Document(
            page_content=content,
            metadata={
                "source": file_path,
                "filename": os.path.basename(file_path),
                "row_count": len(df),
                "column_count": len(df.columns),
                "columns": list(df.columns)
            }
        )
        
        print(f"✅ CSV: {len(df)} rows, {len(df.columns)} columns from {os.path.basename(file_path)}")
        return [doc]
    
    def get_supported_formats(self) -> Dict[str, str]:
        """Get supported file formats"""
        return self.SUPPORTED_FORMATS.copy()
    
    def is_supported(self, file_path: str) -> bool:
        """Check if file format is supported"""
        file_format = self.auto_detect_format(file_path)
        return file_format is not None


# Convenience functions
def create_file_processor() -> FileProcessor:
    """Create file processor instance"""
    return FileProcessor()


def process_file(file_path: str, metadata: Optional[Dict] = None) -> List[Document]:
    """
    Process file with auto-format detection
    
    Args:
        file_path: Path to file
        metadata: Additional metadata
        
    Returns:
        List of Document objects
    """
    processor = create_file_processor()
    return processor.process_file(file_path, metadata)


if __name__ == "__main__":
    # Test file processor
    processor = create_file_processor()
    
    print("📚 Supported Formats:")
    for fmt, desc in processor.get_supported_formats().items():
        print(f"  - {fmt}: {desc}")
