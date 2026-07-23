# loader_improved.py - Chunking cải tiến cho tài liệu học thuật
import os
from langchain_community.document_loaders import PyPDFLoader
from langchain_text_splitters import RecursiveCharacterTextSplitter
from langchain_core.documents import Document
from pptx import Presentation
import re

def load_pdf(file_path):
    """
    Đọc file PDF và trả về danh sách các Document.
    """
    loader = PyPDFLoader(file_path)
    documents = loader.load()
    print(f"Đã đọc {len(documents)} trang từ PDF: {file_path}")
    return documents

def load_ppt(file_path):
    """
    Đọc file PowerPoint (.pptx) và trả về danh sách Document.
    """
    try:
        prs = Presentation(file_path)
        documents = []
        
        for slide_idx, slide in enumerate(prs.slides):
            text_content = []
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text_content.append(shape.text)
            
            content = "\n".join(text_content)
            if content.strip():
                doc = Document(
                    page_content=content,
                    metadata={
                        "source": file_path,
                        "slide": slide_idx + 1,
                        "type": "presentation"
                    }
                )
                documents.append(doc)
        
        print(f"Đã đọc {len(documents)} slide từ PPT: {file_path}")
        return documents
    except Exception as e:
        print(f"⚠️  Không thể đọc {file_path}: {e}")
        return []

def load_docx(file_path):
    """
    Đọc file Word (.docx) bằng python-docx và trả về danh sách Document.
    """
    try:
        from docx import Document as DocxDocument
        doc = DocxDocument(file_path)
        full_text = []
        
        # Đọc tất cả các đoạn văn
        for para in doc.paragraphs:
            if para.text.strip():
                full_text.append(para.text)
        
        # Đọc tất cả các bảng biểu
        for table in doc.tables:
            for row in table.rows:
                row_text = [cell.text.strip() for cell in row.cells if cell.text.strip()]
                if row_text:
                    full_text.append(" | ".join(row_text))
                    
        content = "\n".join(full_text)
        if content.strip():
            doc_node = Document(
                page_content=content,
                metadata={
                    "source": file_path,
                    "type": "docx"
                }
            )
            print(f"Đã đọc file DOCX: {file_path}")
            return [doc_node]
        return []
    except Exception as e:
        print(f"⚠️  Không thể đọc file DOCX {file_path}: {e}")
        return []

def clean_text(text):
    """
    Làm sạch văn bản: loại bỏ khoảng trắng thừa, dòng trống.
    """
    lines = text.split('\n')
    cleaned_lines = [line.strip() for line in lines if line.strip() != '']
    return '\n'.join(cleaned_lines)

def detect_content_type(text):
    """
    Phát hiện loại nội dung để xử lý đặc biệt.
    """
    text_lower = text.lower()
    
    # Phát hiện algorithm/pseudocode
    if any(keyword in text_lower for keyword in ['algorithm', 'procedure', 'input:', 'output:', 'begin', 'end']):
        return 'algorithm'
    
    # Phát hiện table
    if text.count('|') > 5 or 'Table' in text:
        return 'table'
    
    # Phát hiện công thức nhiều
    formula_chars = sum(1 for char in text if char in '∑∫√≤≥∈∀∃⊂⊆∪∩')
    if formula_chars > 10:
        return 'formula_heavy'
    
    # Phát hiện danh sách có số
    if len(re.findall(r'\n\s*\d+\.', text)) > 5:
        return 'numbered_list'
    
    return 'normal'

def chunk_documents_improved(documents, chunk_size=1200, chunk_overlap=250):
    """
    Chunking cải tiến cho tài liệu học thuật.
    
    Cải tiến:
    - Chunk size lớn hơn (1200 vs 1000) để giữ context tốt hơn
    - Overlap lớn hơn (250 vs 200) để tránh mất thông tin
    - Separators chi tiết hơn để ngắt tại ranh giới tự nhiên
    - Xử lý đặc biệt cho algorithms, tables, formulas
    - Lọc bỏ chunks quá ngắn (noise)
    """
    
    # Tạo splitter với cấu hình tối ưu
    text_splitter = RecursiveCharacterTextSplitter(
        chunk_size=chunk_size,
        chunk_overlap=chunk_overlap,
        length_function=len,
        separators=[
            "\n\n\n",    # Ngắt giữa sections lớn (3 dòng trống)
            "\n\n",      # Ngắt giữa đoạn văn (2 dòng trống)
            "\n",        # Ngắt dòng
            ". ",        # Ngắt câu (dấu chấm + space)
            "! ",        # Ngắt câu cảm thán
            "? ",        # Ngắt câu hỏi
            "; ",        # Ngắt mệnh đề
            ", ",        # Ngắt cụm từ
            " ",         # Ngắt từ
            ""           # Fallback: ngắt ký tự
        ]
    )
    
    all_chunks = []
    special_content_count = {
        'algorithm': 0,
        'table': 0,
        'formula_heavy': 0,
        'numbered_list': 0,
        'normal': 0
    }
    
    for doc in documents:
        content = doc.page_content.strip()
        
        # Bỏ qua trang trống hoặc quá ngắn
        if len(content) < 100:
            continue
        
        # Phát hiện loại nội dung
        content_type = detect_content_type(content)
        special_content_count[content_type] += 1
        
        # Xử lý đặc biệt cho algorithms và tables
        if content_type in ['algorithm', 'table']:
            # Nếu nội dung ngắn hơn chunk_size, giữ nguyên
            if len(content) <= chunk_size * 1.5:
                # Thêm metadata đặc biệt
                doc.metadata['content_type'] = content_type
                all_chunks.append(doc)
                continue
        
        # Chunking thông thường
        sub_chunks = text_splitter.split_documents([doc])
        
        # Thêm metadata về content type
        for chunk in sub_chunks:
            if 'content_type' not in chunk.metadata:
                chunk.metadata['content_type'] = content_type
        
        all_chunks.extend(sub_chunks)
    
    # Lọc bỏ chunks quá ngắn (có thể là noise)
    filtered_chunks = [c for c in all_chunks if len(c.page_content.strip()) > 50]
    
    print(f"\n📊 Thống kê chunking:")
    print(f"  - Tổng chunks: {len(filtered_chunks)}")
    print(f"  - Đã lọc bỏ: {len(all_chunks) - len(filtered_chunks)} chunks quá ngắn")
    print(f"\n📝 Phân loại nội dung:")
    for content_type, count in special_content_count.items():
        if count > 0:
            print(f"  - {content_type}: {count} trang")
    
    return filtered_chunks

def chunk_documents(documents, chunk_size=1200, chunk_overlap=250):
    """
    Wrapper function để tương thích với code cũ.
    Sử dụng chunking cải tiến.
    """
    return chunk_documents_improved(documents, chunk_size, chunk_overlap)

def load_all_documents(data_folder="data"):
    """
    Quét thư mục data, đọc tất cả file PDF và PPTX, trả về danh sách chunks.
    """
    all_docs = []
    if not os.path.exists(data_folder):
        print(f"⚠️  Thư mục {data_folder} không tồn tại!")
        return []
    
    print(f"📚 Đang quét thư mục: {data_folder}")
    
    for filename in os.listdir(data_folder):
        file_path = os.path.join(data_folder, filename)
        if filename.endswith(".pdf"):
            docs = load_pdf(file_path)
            all_docs.extend(docs)
        elif filename.endswith(".pptx"):
            docs = load_ppt(file_path)
            all_docs.extend(docs)
        elif filename.endswith(".docx"):
            docs = load_docx(file_path)
            all_docs.extend(docs)
        else:
            continue
    
    if not all_docs:
        print("⚠️  Không tìm thấy file PDF, PPTX hoặc DOCX nào!")
        return []
    
    print(f"\n✅ Đã load {len(all_docs)} documents")
    
    # Làm sạch nội dung
    for doc in all_docs:
        doc.page_content = clean_text(doc.page_content)
    
    # Chia nhỏ thành chunks với chiến lược cải tiến
    chunks = chunk_documents_improved(all_docs)
    
    return chunks

if __name__ == "__main__":
    # Test chunking cải tiến
    print("=" * 70)
    print("TEST CHUNKING CẢI TIẾN")
    print("=" * 70)
    
    chunks = load_all_documents("data")
    
    print(f"\n📊 Kết quả:")
    print(f"  - Tổng số chunks: {len(chunks)}")
    
    # Thống kê độ dài
    lengths = [len(c.page_content) for c in chunks]
    print(f"  - Độ dài trung bình: {sum(lengths)/len(lengths):.0f} chars")
    print(f"  - Độ dài min: {min(lengths)} chars")
    print(f"  - Độ dài max: {max(lengths)} chars")
    
    # Thống kê content types
    content_types = {}
    for chunk in chunks:
        ct = chunk.metadata.get('content_type', 'unknown')
        content_types[ct] = content_types.get(ct, 0) + 1
    
    print(f"\n📝 Phân loại chunks:")
    for ct, count in sorted(content_types.items(), key=lambda x: x[1], reverse=True):
        print(f"  - {ct}: {count} chunks")
    
    # Hiển thị mẫu chunks
    print(f"\n📄 Mẫu chunks:")
    for i in [0, len(chunks)//2, len(chunks)-1]:
        print(f"\n--- Chunk {i+1} ---")
        print(f"Type: {chunks[i].metadata.get('content_type', 'unknown')}")
        print(f"Length: {len(chunks[i].page_content)} chars")
        print(f"Source: {chunks[i].metadata.get('source', 'unknown')}")
        print(f"Content: {chunks[i].page_content[:200]}...")
