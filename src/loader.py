# loader.py - Đọc và xử lý file PDF, PPT
import os
import sys

os.environ['KMP_DUPLICATE_LIB_OK'] = 'True'

if hasattr(sys.stdout, "reconfigure"):
    try:
        sys.stdout.reconfigure(encoding="utf-8", errors="replace")
    except TypeError:
        sys.stdout.reconfigure(errors="replace")
    except Exception:
        pass
from langchain_community.document_loaders import PyMuPDFLoader
from langchain_text_splitters import RecursiveCharacterTextSplitter
from langchain_core.documents import Document
from pptx import Presentation

def _extract_docx_xml_text(file_path: str) -> str:
    """Fallback DOCX reader for text boxes, headers/footers, and unusual Word XML."""
    import re
    import zipfile
    import xml.etree.ElementTree as ET

    xml_targets = (
        "word/document.xml",
        "word/footnotes.xml",
        "word/endnotes.xml",
        "word/comments.xml",
        "word/glossary/document.xml",
    )

    def local_name(tag: str) -> str:
        return tag.rsplit("}", 1)[-1] if "}" in tag else tag

    def paragraph_text(paragraph) -> str:
        parts = []
        for node in paragraph.iter():
            name = local_name(node.tag)
            if name == "t" and node.text:
                parts.append(node.text)
            elif name == "tab":
                parts.append("\t")
            elif name in {"br", "cr"}:
                parts.append("\n")
        return re.sub(r"[ \t]+", " ", "".join(parts)).strip()

    lines = []
    with zipfile.ZipFile(file_path) as archive:
        names = set(archive.namelist())
        selected = [name for name in xml_targets if name in names]
        selected.extend(
            name for name in names
            if re.match(r"word/(header|footer)\d+\.xml$", name)
        )

        for name in selected:
            try:
                root = ET.fromstring(archive.read(name))
            except Exception:
                continue
            for paragraph in root.iter():
                if local_name(paragraph.tag) != "p":
                    continue
                text = paragraph_text(paragraph)
                if text:
                    lines.append(text)

    deduped = []
    previous = ""
    for line in lines:
        if line != previous:
            deduped.append(line)
        previous = line
    return "\n".join(deduped).strip()

def load_pdf(file_path):
    """
    Đọc file PDF và trả về danh sách các Document.
    Mỗi Document chứa nội dung 1 trang và metadata chuẩn hóa.
    """
    import hashlib
    from datetime import datetime
    from config import EMBEDDING_MODEL_NAME

    loader = PyMuPDFLoader(file_path)
    documents = loader.load()

    filename = os.path.basename(file_path)
    now = datetime.utcnow().isoformat()

    for doc in documents:
        page = doc.metadata.get("page", 0)
        chunk_id = hashlib.md5(f"{filename}:p{page}:{doc.page_content[:80]}".encode()).hexdigest()[:16]
        doc.metadata.update({
            "chunk_id": chunk_id,
            "source_file": filename,
            "source": file_path,
            "page_number": page + 1,
            "page": page + 1,
            "section": f"Page {page + 1}",
            "file_type": "pdf",
            "created_at": now,
            "embedding_model": EMBEDDING_MODEL_NAME,
        })

    print(f"Đã đọc {len(documents)} trang từ PDF: {file_path}")
    return documents

def load_ppt(file_path):
    """
    Đọc file PowerPoint (.pptx) và trả về danh sách Document với metadata chuẩn hóa.
    """
    import hashlib
    from datetime import datetime
    from config import EMBEDDING_MODEL_NAME

    try:
        prs = Presentation(file_path)
        documents = []
        filename = os.path.basename(file_path)
        now = datetime.utcnow().isoformat()

        for slide_idx, slide in enumerate(prs.slides):
            text_content = []
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text_content.append(shape.text)

            content = "\n".join(text_content)
            if content.strip():
                slide_num = slide_idx + 1
                chunk_id = hashlib.md5(f"{filename}:s{slide_num}:{content[:80]}".encode()).hexdigest()[:16]
                doc = Document(
                    page_content=content,
                    metadata={
                        "chunk_id": chunk_id,
                        "source_file": filename,
                        "source": file_path,
                        "page_number": slide_num,
                        "slide": slide_num,
                        "section": f"Slide {slide_num}",
                        "file_type": "pptx",
                        "created_at": now,
                        "embedding_model": EMBEDDING_MODEL_NAME,
                    }
                )
                documents.append(doc)

        print(f"Đã đọc {len(documents)} slide từ PPT: {file_path}")
        return documents
    except Exception as e:
        print(f"⚠️  Không thể đọc {file_path}: {e}")
        return []

def load_docx(file_path):
    """
    Đọc file Word (.docx) — parser nâng cấp.

    Xử lý đầy đủ:
      - Heading (Word style + Markdown ##)
      - Paragraph thông thường
      - Bullet list (giữ ký tự •/-)
      - Table → chuyển thành văn bản có ngữ nghĩa
      - Nested table (table trong cell)

    Table rendering:
      Thay vì "Tiêu chí | Điểm" (pipe-join vô nghĩa)
      → "Tiêu chí: Điểm" hoặc "Hình thức đánh giá: Kiểm tra lý thuyết | Tỷ lệ: 25%"
      Hàng tiêu đề được nhận dạng và dùng để đặt nhãn cho các ô dữ liệu.
    """
    if os.path.basename(file_path).startswith("~$"):
        return []

    try:
        from docx import Document as DocxDocument
        import re
        import hashlib
        from datetime import datetime
        from config import EMBEDDING_MODEL_NAME

        doc = DocxDocument(file_path)
        filename = os.path.basename(file_path)
        now = datetime.utcnow().isoformat()
        documents = []

        current_section_title = ""
        current_section_lines: list = []
        section_index = 1

        _md_heading = re.compile(r'^#{2,4}\s+\S')

        # ── Helpers ──────────────────────────────────────────────────────────

        def _is_heading(para) -> bool:
            """Word Heading style hoặc Markdown ## heading."""
            if para.style.name.startswith("Heading"):
                return True
            text = para.text.strip()
            if _md_heading.match(text):
                if len(text.lstrip('#').strip()) >= 8:
                    return True
            return False

        def _is_bullet(para) -> bool:
            """Paragraph là bullet/list item."""
            style = para.style.name.lower()
            if "list" in style or "bullet" in style:
                return True
            # numFmt hoặc numId trong XML
            if para._element.find('.//{http://schemas.openxmlformats.org/wordprocessingml/2006/main}numPr') is not None:
                return True
            return False

        def _cell_text(cell) -> str:
            """Lấy toàn bộ text trong cell, kể cả nested table."""
            parts = []
            for para in cell.paragraphs:
                t = para.text.strip()
                if t:
                    parts.append(t)
            # nested table: đệ quy lấy text
            for nested_table in cell.tables:
                nested_lines = _render_table(nested_table, header_row=None)
                parts.extend(nested_lines)
            return " ".join(parts)

        def _detect_header_row(table):
            """
            Phát hiện hàng tiêu đề của bảng.
            Hàng header thường là hàng đầu tiên (hoặc 2 hàng đầu nếu merged).
            Trả về list of header strings, hoặc None nếu không xác định được.
            """
            if not table.rows:
                return None
            first_row_cells = [_cell_text(c).strip() for c in table.rows[0].cells]
            # Nếu hàng đầu có ít nhất 2 cell không rỗng → coi là header
            non_empty = [c for c in first_row_cells if c]
            if len(non_empty) >= 2:
                return first_row_cells
            return None

        def _render_table(table, header_row=None) -> list:
            """
            Chuyển bảng thành list of strings có ngữ nghĩa.

            Chiến lược:
            1. Nếu có header row: mỗi cell dữ liệu = "Header: Giá trị"
            2. Nếu bảng 2 cột (key-value): "Key: Value"
            3. Fallback: join các cell bằng " | " nhưng lọc bỏ cell rỗng
            """
            if not table.rows:
                return []

            lines = []
            detected_header = header_row if header_row is not None else _detect_header_row(table)

            # Xác định hàng bắt đầu data (bỏ qua header rows)
            data_start = 0
            if detected_header:
                data_start = 1
                # Kiểm tra nếu hàng 2 cũng là header lặp lại (merged header 2 tầng)
                if len(table.rows) > 2:
                    row2_cells = [_cell_text(c).strip() for c in table.rows[1].cells]
                    seen2 = set()
                    unique2 = []
                    for c in row2_cells:
                        if c not in seen2:
                            unique2.append(c)
                            seen2.add(c)
                    # Nếu hàng 2 là tập con của header → cũng là header
                    non_empty2 = set(c for c in unique2 if c)
                    non_empty_h = set(c for c in detected_header if c)
                    if non_empty2 and non_empty2.issubset(non_empty_h):
                        data_start = 2

            num_cols = len(table.columns)

            for row_idx, row in enumerate(table.rows):
                if row_idx < data_start:
                    continue  # bỏ header

                cells = [_cell_text(c).strip() for c in row.cells]

                # Loại bỏ các cell trùng lặp do merged cells
                unique_cells = []
                seen = set()
                for c in cells:
                    if c not in seen:
                        unique_cells.append(c)
                        seen.add(c)

                non_empty = [c for c in unique_cells if c]
                if not non_empty:
                    continue

                # Nếu chỉ có 1 cell có nội dung → ghi trực tiếp (không prefix header)
                if len(non_empty) == 1:
                    lines.append(non_empty[0])
                    continue

                # Chiến lược 1: có header → "Header_i: cell_i"
                if detected_header and len(unique_cells) >= 2:
                    parts = []
                    for col_i, cell_val in enumerate(unique_cells):
                        if not cell_val:
                            continue
                        header_label = ""
                        if col_i < len(detected_header):
                            hl = detected_header[col_i].strip()
                            # Bỏ header trống, số thứ tự đơn giản, hoặc trùng với value
                            if hl and not re.match(r'^[\d.]+$', hl) and hl != cell_val:
                                header_label = hl
                        if header_label:
                            parts.append(f"{header_label}: {cell_val}")
                        else:
                            parts.append(cell_val)
                    if parts:
                        lines.append(" | ".join(parts))

                # Chiến lược 2: bảng 2 cột key-value (không có detected_header)
                elif num_cols == 2 and len(non_empty) == 2:
                    lines.append(f"{non_empty[0]}: {non_empty[1]}")

                # Fallback
                else:
                    lines.append(" | ".join(non_empty))

            return lines

        # ── Flush section ─────────────────────────────────────────────────────

        def flush_section():
            nonlocal section_index
            content = "\n".join(current_section_lines).strip()
            if not content:
                return
            chunk_id = hashlib.md5(
                f"{filename}:sec{section_index}:{content[:80]}".encode()
            ).hexdigest()[:16]
            documents.append(Document(
                page_content=content,
                metadata={
                    "chunk_id": chunk_id,
                    "source_file": filename,
                    "source": file_path,
                    "page_number": section_index,
                    "page": section_index,
                    "section": current_section_title or f"Section {section_index}",
                    "file_type": "docx",
                    "type": "docx",
                    "created_at": now,
                    "embedding_model": EMBEDDING_MODEL_NAME,
                }
            ))
            section_index += 1

        # ── Main traversal ────────────────────────────────────────────────────
        # Duyệt body elements theo đúng thứ tự (paragraph xen kẽ với table)

        body = doc.element.body
        para_idx = 0
        table_idx = 0
        para_list = doc.paragraphs
        table_list = doc.tables

        for child in body:
            tag = child.tag.split("}")[-1] if "}" in child.tag else child.tag

            if tag == "p":
                if para_idx >= len(para_list):
                    para_idx += 1
                    continue
                para = para_list[para_idx]
                para_idx += 1
                text = para.text.strip()
                if not text:
                    continue

                if _is_heading(para):
                    flush_section()
                    current_section_title = text
                    current_section_lines = [text]
                elif _is_bullet(para):
                    current_section_lines.append(f"• {text}")
                else:
                    current_section_lines.append(text)

            elif tag == "tbl":
                if table_idx >= len(table_list):
                    table_idx += 1
                    continue
                table = table_list[table_idx]
                table_idx += 1

                table_lines = _render_table(table)
                if table_lines:
                    current_section_lines.extend(table_lines)

        flush_section()

        # Fallback: nếu không có section nào (file không có heading)
        if not documents:
            all_lines = []
            for para in doc.paragraphs:
                t = para.text.strip()
                if t:
                    if _is_bullet(para):
                        all_lines.append(f"• {t}")
                    else:
                        all_lines.append(t)
            for table in doc.tables:
                all_lines.extend(_render_table(table))
            content = "\n".join(all_lines).strip()
            if content:
                chunk_id = hashlib.md5(f"{filename}:fallback:{content[:80]}".encode()).hexdigest()[:16]
                documents.append(Document(
                    page_content=content,
                    metadata={
                        "chunk_id": chunk_id,
                        "source_file": filename,
                        "source": file_path,
                        "page_number": 1,
                        "page": 1,
                        "section": "Full document",
                        "file_type": "docx",
                        "type": "docx",
                        "created_at": now,
                        "embedding_model": EMBEDDING_MODEL_NAME,
                    }
                ))

        if not documents:
            content = _extract_docx_xml_text(file_path)
            if content:
                chunk_id = hashlib.md5(f"{filename}:xmlfallback:{content[:80]}".encode()).hexdigest()[:16]
                documents.append(Document(
                    page_content=content,
                    metadata={
                        "chunk_id": chunk_id,
                        "source_file": filename,
                        "source": file_path,
                        "page_number": 1,
                        "page": 1,
                        "section": "Full document",
                        "file_type": "docx",
                        "type": "docx",
                        "created_at": now,
                        "embedding_model": EMBEDDING_MODEL_NAME,
                        "extraction_method": "docx_xml_fallback",
                    }
                ))

        print(f"Đã đọc {len(documents)} section từ DOCX: {file_path}")
        return documents

    except Exception as e:
        try:
            import hashlib
            from datetime import datetime
            from config import EMBEDDING_MODEL_NAME

            filename = os.path.basename(file_path)
            now = datetime.utcnow().isoformat()
            content = _extract_docx_xml_text(file_path)
            if content:
                chunk_id = hashlib.md5(f"{filename}:xmlrecover:{content[:80]}".encode()).hexdigest()[:16]
                print(f"⚠️  DOCX parser chính lỗi, dùng XML fallback cho {file_path}: {e}")
                return [Document(
                    page_content=content,
                    metadata={
                        "chunk_id": chunk_id,
                        "source_file": filename,
                        "source": file_path,
                        "page_number": 1,
                        "page": 1,
                        "section": "Full document",
                        "file_type": "docx",
                        "type": "docx",
                        "created_at": now,
                        "embedding_model": EMBEDDING_MODEL_NAME,
                        "extraction_method": "docx_xml_recovery",
                    }
                )]
        except Exception as fallback_error:
            print(f"⚠️  DOCX XML fallback cũng thất bại {file_path}: {fallback_error}")

        print(f"⚠️  Không thể đọc file DOCX {file_path}: {e}")
        import traceback
        traceback.print_exc()
        return []

def clean_text(text):
    """
    Làm sạch văn bản cơ bản: loại bỏ khoảng trắng thừa, dòng trống.
    (Có thể thêm xử lý khác nếu cần)
    """
    # Xóa khoảng trắng đầu cuối dòng
    lines = text.split('\n')
    cleaned_lines = [line.strip() for line in lines if line.strip() != '']
    return '\n'.join(cleaned_lines)

def chunk_documents(documents, chunk_size=1000, chunk_overlap=200, use_semantic=False):
    """
    Chia nhỏ danh sách Document thành các chunks.

    Chiến lược (theo thứ tự ưu tiên):
      Level 1 — Heading-based (DOCX): bảo toàn heading + nội dung
      Level 2 — Section-based (PDF/PPTX): mỗi trang/slide là 1 chunk
      Level 3 — Sentence-split (fallback): khi chunk vẫn quá lớn

    use_semantic=True giữ tương thích ngược nhưng không còn dùng
    SemanticChunker vì nó load thêm embedding model (chậm, tốn RAM).
    """
    cleaned_docs = [doc for doc in documents if doc.page_content and doc.page_content.strip()]
    if not cleaned_docs:
        return []

    from chunking.academic_chunker import smart_chunk
    chunks = smart_chunk(cleaned_docs)
    if chunks:
        return chunks

    print("⚠️ smart_chunk returned 0 chunks; falling back to RecursiveCharacterTextSplitter")
    splitter = RecursiveCharacterTextSplitter(
        chunk_size=chunk_size,
        chunk_overlap=chunk_overlap,
        separators=["\n\n", "\n", ". ", " ", ""],
    )
    fallback_chunks = splitter.split_documents(cleaned_docs)
    return [chunk for chunk in fallback_chunks if chunk.page_content.strip()]

def load_all_documents(data_folder="data", use_semantic_chunking=True):
    """
    Quét thư mục data, đọc tất cả file PDF và PPTX, trả về danh sách chunks.
    
    Args:
        data_folder: Thư mục chứa dữ liệu
        use_semantic_chunking: Sử dụng semantic chunking hay không
    
    Returns:
        Danh sách chunks
    
    Lưu ý: Chỉ hỗ trợ .pptx, không hỗ trợ .ppt cũ
    """
    all_docs = []
    if not os.path.exists(data_folder):
        print(f"⚠️  Thư mục {data_folder} không tồn tại!")
        return []
    
    for filename in os.listdir(data_folder):
        file_path = os.path.join(data_folder, filename)
        if filename.endswith(".pdf"):
            docs = load_pdf(file_path)
            all_docs.extend(docs)
        elif filename.endswith(".pptx"):  # Chỉ .pptx, không .ppt
            docs = load_ppt(file_path)
            all_docs.extend(docs)
        elif filename.endswith(".docx"):
            docs = load_docx(file_path)
            all_docs.extend(docs)
        else:
            continue
    
    if not all_docs:
        print("⚠️  Không tìm thấy file PDF, PPTX hoặc DOCX nào!")
        return []
    
    # Làm sạch nội dung từng document
    for doc in all_docs:
        doc.page_content = clean_text(doc.page_content)
    
    # Chia nhỏ thành chunks dùng chiến lược 3-level mới
    raw_chunks = chunk_documents(all_docs)
    
    chunks = []
    for chunk in raw_chunks:
        filename = chunk.metadata.get("source_file", "")
        section = chunk.metadata.get("section", "")
        content = chunk.page_content
        
        # Xử lý đặc biệt cho thuattoanpython.docx:
        # Bỏ qua phần "I. TỪ VIẾT TẮT VÀ THUẬT NGỮ"
        if filename == "thuattoanpython.docx":
            if "I. TỪ VIẾT TẮT" in section or "I. TỪ VIẾT TẮT" in content or "Viết tắt | Tiếng Anh" in content:
                continue
        
        # A. Nhóm Syllabus
        if filename in ["220269_ Khai pha du lieu - AI.docx", "220269_ Khai pha du lieu - CNTT.docx"]:
            chunk.metadata["doc_type"] = "syllabus"
            chunk.metadata["lesson"] = "all"
        # B. Nhóm Theory
        elif filename == "01Intro.pptx":
            chunk.metadata["doc_type"] = "theory"
            chunk.metadata["lesson"] = "bai_1"
        elif filename in ["02Data.pptx", "03Preprocessing.pptx", "04OLAP.pptx", "05CubeTech.pptx"]:
            chunk.metadata["doc_type"] = "theory"
            chunk.metadata["lesson"] = "bai_2"
        elif filename in ["06FPBasic.pptx", "07FPAdvanced.pptx"]:
            chunk.metadata["doc_type"] = "theory"
            chunk.metadata["lesson"] = "bai_3"
        elif filename in ["08ClassBasic.pptx", "09ClassAdvanced.pptx"]:
            chunk.metadata["doc_type"] = "theory"
            chunk.metadata["lesson"] = "bai_4"
        elif filename in ["10ClusBasic.pptx", "11ClusAdvanced.pptx"]:
            chunk.metadata["doc_type"] = "theory"
            chunk.metadata["lesson"] = "bai_5"
        elif filename in ["DM3.pdf", "Data.Mining.Concepts.and.Techniques.2nd.Ed-1558609016.pdf"]:
            chunk.metadata["doc_type"] = "theory"
            chunk.metadata["lesson"] = "all"
        elif filename.endswith(".pptx"):
            # Các bài giảng khác (ví dụ: OLAP, CubeTech, Outlier, Trend)
            chunk.metadata["doc_type"] = "theory"
            chunk.metadata["lesson"] = "all"
        # C. Nhóm Code
        elif filename == "thuattoanpython.docx":
            chunk.metadata["doc_type"] = "code"
            chunk.metadata["lesson"] = "all"
        else:
            # Fallback
            chunk.metadata["doc_type"] = "theory"
            chunk.metadata["lesson"] = "all"
            
        chunks.append(chunk)
        
    return chunks

if __name__ == "__main__":
    # Test thử: chạy file này để xem có đọc được không
    chunks = load_all_documents("data")
    print(f"Tổng số chunks: {len(chunks)}")
    print("Nội dung chunk đầu tiên:")
    print(chunks[0].page_content[:500])  # In 500 ký tự đầu
